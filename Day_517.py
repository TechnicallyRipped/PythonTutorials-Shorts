

# pip install python-pptx
from pptx import Presentation

prs = Presentation('text.pptx')

slides = prs.slides 

for i,slide in enumerate(slides):
    for shape in slide.shapes:
        print(f'{i} text: {shape.text}')


