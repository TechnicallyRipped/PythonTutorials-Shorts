

# pip install python-pptx
from pptx import Presentation

prs = Presentation()

layout = prs.slide_layouts[0]

slide = prs.slides.add_slide(layout)

title = slide.shapes.title
title.text = 'TechnicallyRipped'

prs.save('automated_pp.pptx')