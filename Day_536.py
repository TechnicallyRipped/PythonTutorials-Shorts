
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

left = Inches(2) # Left position
top = Inches(2) # Top position
width = Inches(6) # Width of whole table
height = Inches(1) # Height of each row

table = slide.shapes.add_table(3,4,left,top,
                               width,height).table
table.cell(0,1).text = 'TEST'

prs.save("table.pptx")
