


from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

pic_path = 'pic_to_add.png'

left = Inches(1)
top = Inches(2)

slide.shapes.add_picture(pic_path, left, top)

prs.save('pres_with_pic.pptx')



























# import win32com.client

# powerpoint_file = 'text.pptx'
# pdf_file = 'text_powerpoint.pdf'

# powerpoint = win32com.client.Dispatch("PowerPoint.Application")

# presentation = powerpoint.Presentations.Open(powerpoint_file, WithWindow=False)

# # presentation.SaveAs(pdf_file, 32)
# # presentation.Close()
# # powerpoint.Quit()

# print('Done converting')



















# pip install pypandoc
# import pypandoc

# pypandoc.download_pandoc()

# powerpoint_file = 'text.pptx'
# pdf_file = 'text_powerpoint.pdf'
# pypandoc.convert_file(powerpoint_file,'pdf',
#                       outputfile=pdf_file)

# print('Done converting!')