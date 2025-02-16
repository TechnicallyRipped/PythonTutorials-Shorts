

from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

left = Inches(1)
top = Inches(1)
width = Inches(6)
height = Inches(3)

slide.shapes.add_picture("sub.jpg", left, top, width, height)

prs.save("output_presentation.pptx")
