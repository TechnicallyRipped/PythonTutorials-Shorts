
# pip install python-pptx
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
blank_slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide)

# Height: 7.5 inches \\ Width: 10 inches
left = Inches(1)
top = Inches(1) 
width = Inches(4)
height = Inches(1)  

tb = slide.shapes.add_textbox(left, top, 
                              width, height)

tb_text = tb.text_frame
tb_text.text = "TechnicallyRipped"

prs.save('text.pptx')
